"""Generate euAIKG technical presentation (6 slides)."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DOCS = Path(__file__).parent
OUT = DOCS / "euAIKG_technical.pptx"

# Slide dimensions: 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x1E, 0x88, 0xE5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, left, top, width, height, size=32, color=DARK, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tf


def add_body(slide, lines, left, top, width, height, size=14, color=DARK, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(size * spacing * 0.4)
    return tf


def add_subtitle(slide, text, left, top, width, height, size=16, color=ACCENT):
    return add_title(slide, text, left, top, width, height, size=size, color=color, bold=True)


def add_image_centered(slide, img_path, top, max_w, max_h):
    """Add image centered horizontally, fitting within max_w x max_h."""
    from PIL import Image
    img = Image.open(img_path)
    iw, ih = img.size
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    left = (SLIDE_W - Emu(w * 914400 // 96)) // 2  # center
    slide.shapes.add_picture(str(img_path), left, top, Emu(w * 914400 // 96), Emu(h * 914400 // 96))


def make_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # blank layout

    # ================================================================
    # SLIDE 1 — Architecture
    # ================================================================
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, WHITE)
    add_title(s1, "시스템 아키텍처", Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), size=30)
    add_body(s1, [
        "EU AI Act 원문 -> LLM 기반 그래프 추출 -> Neo4j 저장 -> 웹 대시보드 시각화",
        "",
        "- 백엔드: Neo4j (그래프 DB) + vLLM / Gemini API (계산엔진)",
        "- 프론트엔드: Flask + Cytoscape.js (단일 페이지 대시보드)",
        "- 재개 가능한 파이프라인, CLI / 웹 이중 실행 모드 지원",
    ], Inches(0.8), Inches(1.1), Inches(5.5), Inches(2.0), size=13)

    img1 = DOCS / "slide1_architecture.png"
    if img1.exists():
        s1.shapes.add_picture(str(img1), Inches(0.5), Inches(3.3), Inches(12.3))

    # ================================================================
    # SLIDE 2 — Principles
    # ================================================================
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, WHITE)
    add_title(s2, "설계 원칙", Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), size=30)

    principles = [
        ("이중 LLM 견고성",
         "로컬 vLLM으로 처리량 확보, 실패 청크만 Gemini API로 재시도\n"
         "-> 비용 효율성 + 추출 완성도 동시 달성"),
        ("그래프 네이티브 저장",
         "엔티티 간 관계를 Neo4j에서 1차 시민(first-class)으로 표현\n"
         "GDS 플러그인을 통한 네이티브 그래프 알고리즘(KNN, WCC) 적용"),
        ("LLM 기반 엔티티 해소",
         "임베딩 유사도 + 텍스트 거리로 후보군 축소\n"
         "Gemini 구조화 출력으로 최종 병합 판정"),
        ("재개 가능한 파이프라인",
         "단계별 체크포인트(pickle 존재, 노드 수, WCC 속성) 확인\n"
         "중단 후 재처리 없이 이어서 실행 가능"),
    ]

    for i, (title, desc) in enumerate(principles):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 6.2)
        top = Inches(1.3 + row * 2.8)
        # Card background
        shape = s2.shapes.add_shape(
            1, left, top, Inches(5.6), Inches(2.3)  # 1 = rectangle
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF6, 0xF8)
        shape.line.fill.background()
        shape.shadow.inherit = False

        add_subtitle(s2, f"{i+1}. {title}", left + Inches(0.3), top + Inches(0.2),
                     Inches(5.0), Inches(0.5), size=16)
        add_body(s2, desc.split("\n"), left + Inches(0.3), top + Inches(0.8),
                 Inches(5.0), Inches(1.3), size=12, color=GRAY)

    # ================================================================
    # SLIDE 3 — Methodology 1
    # ================================================================
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, WHITE)
    add_title(s3, "방법론 — 그래프 추출 및 적재", Inches(0.8), Inches(0.3),
              Inches(11), Inches(0.7), size=28)

    add_body(s3, [
        "그래프 추출",
        "  - LLMGraphTransformer로 텍스트 청크에서 엔티티·관계 자동 추출",
        "  - 4스레드 병렬 처리, 청크당 300초 타임아웃",
        "  - vLLM (Qwen3-14B-AWQ) 1차 추출 + Gemini API 2차 폴백",
        "",
        "적재",
        "  - 역직렬화된 GraphDocument를 Neo4j에 일괄 적재",
        "  - __Entity__ 공통 레이블 부여, MENTIONS 원문 참조 관계 생성",
    ], Inches(0.8), Inches(1.1), Inches(5.0), Inches(2.5), size=12)

    img3 = DOCS / "slide3_extraction.png"
    if img3.exists():
        s3.shapes.add_picture(str(img3), Inches(0.3), Inches(4.0), Inches(12.5))

    # ================================================================
    # SLIDE 4 — Methodology 2
    # ================================================================
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4, WHITE)
    add_title(s4, "방법론 — 커뮤니티 탐지 및 엔티티 해소", Inches(0.8), Inches(0.3),
              Inches(11), Inches(0.7), size=28)

    add_body(s4, [
        "1단계: 임베딩",
        "  - multilingual-e5-large -> 엔티티별 1024차원 벡터 생성",
        "",
        "2단계: 커뮤니티 탐지 (GDS)",
        "  - KNN: 코사인 유사도 >= 0.95 -> SIMILAR 관계 생성",
        "  - WCC: 연결 컴포넌트 -> 커뮤니티 ID 부여",
        "",
        "3단계: 엔티티 해소",
        "  - 후보 필터링: APOC 텍스트 편집거리 <= 3",
        "  - Gemini 구조화 출력(Pydantic 스키마)으로 병합 판정",
        "  - APOC mergeNodes로 중복 노드 병합",
    ], Inches(0.8), Inches(1.1), Inches(5.5), Inches(3.0), size=12)

    img4 = DOCS / "slide4_community.png"
    if img4.exists():
        s4.shapes.add_picture(str(img4), Inches(0.3), Inches(4.5), Inches(12.5))

    # ================================================================
    # SLIDE 5 — Results 1 (Dashboard)
    # ================================================================
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5, WHITE)
    add_title(s5, "결과 — 인터랙티브 대시보드", Inches(0.8), Inches(0.3),
              Inches(11), Inches(0.7), size=30)

    add_body(s5, [
        "- 파이프라인 제어 + 그래프 시각화 + 실시간 로그를 단일 대시보드에 통합",
        "- SSE 기반 실시간 로그 스트리밍, 파이프라인 완료 시 그래프 자동 갱신",
        "- 5종 레이아웃 알고리즘 전환, 팬/줌 인터랙션 지원",
    ], Inches(0.8), Inches(1.0), Inches(11), Inches(1.2), size=13)

    img_overview = DOCS / "screenshot_graph_overview.png"
    img_detail = DOCS / "screenshot_graph_detail.png"
    if img_overview.exists():
        s5.shapes.add_picture(str(img_overview), Inches(0.5), Inches(2.5), Inches(6.0))
    if img_detail.exists():
        s5.shapes.add_picture(str(img_detail), Inches(6.8), Inches(2.5), Inches(6.0))

    # Caption
    add_body(s5, [
        "좌: 전체 지식 그래프 네트워크 뷰                    우: 엔티티 간 관계 상세 뷰"
    ], Inches(0.8), Inches(6.8), Inches(11), Inches(0.5), size=11, color=GRAY)

    # ================================================================
    # SLIDE 6 — Results 2 (Stats)
    # ================================================================
    s6 = prs.slides.add_slide(blank)
    set_slide_bg(s6, WHITE)
    add_title(s6, "결과 — 파이프라인 산출물", Inches(0.8), Inches(0.3),
              Inches(11), Inches(0.7), size=30)

    # Table
    rows, cols = 7, 2
    tbl_shape = s6.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(7), Inches(3.5))
    tbl = tbl_shape.table

    headers = ["항목", "값"]
    data = [
        ("입력 문서", "EU AI Act 원문 (텍스트)"),
        ("청킹", "Qwen3 토크나이저, 350토큰 / 75 오버랩"),
        ("추출 전략", "vLLM 1차 추출 + Gemini 2차 폴백"),
        ("최종 노드 수", "475개 (__Entity__)"),
        ("최종 엣지 수", "500개 (관계)"),
        ("엔티티 해소", "KNN+WCC 커뮤니티 탐지 -> Gemini 판정 -> APOC 병합"),
    ]

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK

    for ri, (k, v) in enumerate(data, start=1):
        for ci, val in enumerate([k, v]):
            cell = tbl.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF6, 0xF8) if ri % 2 == 0 else WHITE

    # Key achievements
    add_subtitle(s6, "주요 성과", Inches(1.5), Inches(5.3), Inches(7), Inches(0.4), size=16)
    add_body(s6, [
        "- EU AI Act를 구조화된 지식 그래프로 자동 변환",
        "- 이중 LLM 전략으로 단일 모델 대비 추출 실패율 감소",
        "- 커뮤니티 기반 엔티티 해소로 중복 노드 제거, 그래프 품질 향상",
    ], Inches(1.5), Inches(5.8), Inches(9), Inches(1.5), size=12)

    # ── Save ──
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    make_pptx()
